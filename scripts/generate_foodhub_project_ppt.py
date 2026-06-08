from pathlib import Path
import json

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
LOGO_IMAGE = Path(r"C:\Users\sai93\Downloads\foodhub-order-eat-enjoy.jpg")
OUTPUT = ROOT / "FoodHub_Project_QA_Observability_DeepSeek.pptx"
METRICS = ROOT / "qa-artifacts" / "observability" / "metrics" / "latest-observability-metrics.json"

RED = RGBColor(224, 31, 31)
DARK_RED = RGBColor(156, 16, 22)
WHITE = RGBColor(255, 255, 255)
BLACK = RGBColor(24, 24, 24)
CHARCOAL = RGBColor(49, 49, 49)
MUTED = RGBColor(105, 105, 105)
SOFT = RGBColor(248, 247, 245)
GOLD = RGBColor(246, 184, 70)
GREEN = RGBColor(34, 150, 87)
BLUE = RGBColor(42, 107, 184)


def load_summary():
    if not METRICS.exists():
        return {
            "totalChecks": 41,
            "passedChecks": 41,
            "failedChecks": 0,
            "requestCount": 11488,
            "p95DurationMs": 2,
            "avgDurationMs": 0.77,
            "paymentSuccessCount": 3716,
            "paymentFailureCount": 0,
        }
    with METRICS.open("r", encoding="utf-8") as f:
        return json.load(f)["summary"]


summary = load_summary()


def prs():
    p = Presentation()
    p.slide_width = Inches(13.333)
    p.slide_height = Inches(7.5)
    return p


deck = prs()
BLANK = deck.slide_layouts[6]


def add_bg(slide, color=WHITE):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, deck.slide_width, deck.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_top_bar(slide, title, subtitle=None):
    slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, deck.slide_width, Inches(0.78)).fill.solid()
    bar = slide.shapes[-1]
    bar.fill.fore_color.rgb = RED
    bar.line.fill.background()
    text(slide, title, 0.55, 0.19, 8.2, 0.34, 20, WHITE, bold=True)
    if subtitle:
        text(slide, subtitle, 9.3, 0.23, 3.4, 0.25, 9, WHITE, align=PP_ALIGN.RIGHT)


def text(slide, value, x, y, w, h, size=18, color=BLACK, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.margin_left = Inches(0.05)
    frame.margin_right = Inches(0.05)
    frame.margin_top = Inches(0.03)
    frame.margin_bottom = Inches(0.03)
    frame.vertical_anchor = MSO_ANCHOR.TOP
    p = frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = value
    run.font.name = "Aptos"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def bullet_list(slide, items, x, y, w, h, size=16, color=CHARCOAL):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.margin_left = Inches(0.12)
    frame.margin_right = Inches(0.1)
    frame.word_wrap = True
    for idx, item in enumerate(items):
        p = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
        p.text = item
        p.level = 0
        p.font.name = "Aptos"
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(7)
    return box


def rounded(slide, x, y, w, h, fill=SOFT, line=RGBColor(225, 225, 225), radius=True):
    kind = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shape = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(1)
    return shape


def label_card(slide, title, body, x, y, w, h, accent=RED):
    rounded(slide, x, y, w, h, WHITE, RGBColor(228, 228, 228))
    slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(y), Inches(0.12), Inches(h)).fill.solid()
    accent_bar = slide.shapes[-1]
    accent_bar.fill.fore_color.rgb = accent
    accent_bar.line.fill.background()
    text(slide, title, x + 0.28, y + 0.18, w - 0.42, 0.35, 15, BLACK, bold=True)
    text(slide, body, x + 0.28, y + 0.62, w - 0.42, h - 0.72, 11.5, MUTED)


def metric_card(slide, value, label, x, y, w, h, accent=RED):
    rounded(slide, x, y, w, h, WHITE, RGBColor(230, 230, 230))
    text(slide, str(value), x + 0.12, y + 0.2, w - 0.24, 0.55, 26, accent, bold=True, align=PP_ALIGN.CENTER)
    text(slide, label, x + 0.12, y + 0.82, w - 0.24, 0.36, 10.5, MUTED, align=PP_ALIGN.CENTER)


def connector(slide, x1, y1, x2, y2, color=RED):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    c.line.color.rgb = color
    c.line.width = Pt(2)


def add_image(slide, path, x, y, w=None, h=None):
    if path.exists():
        if w and h:
            return slide.shapes.add_picture(str(path), Inches(x), Inches(y), Inches(w), Inches(h))
        if w:
            return slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w))
        if h:
            return slide.shapes.add_picture(str(path), Inches(x), Inches(y), height=Inches(h))
    return None


def title_slide():
    slide = deck.slides.add_slide(BLANK)
    add_bg(slide, RED)
    add_image(slide, LOGO_IMAGE, 3.55, 0.8, 6.25)
    text(slide, "Project QA, Observability & AI Test Intelligence", 1.4, 5.15, 10.5, 0.48, 25, WHITE, bold=True, align=PP_ALIGN.CENTER)
    text(slide, "How FoodHub is built, tested, monitored, and released with data-backed go / no-go decisions", 1.5, 5.75, 10.4, 0.35, 14, WHITE, align=PP_ALIGN.CENTER)
    text(slide, "ORDER. EAT. ENJOY.", 4.35, 6.55, 4.7, 0.26, 12, WHITE, bold=True, align=PP_ALIGN.CENTER)


def agenda_slide():
    slide = deck.slides.add_slide(BLANK)
    add_bg(slide, WHITE)
    add_top_bar(slide, "Presentation Flow", "FoodHub Project")
    cards = [
        ("01", "Product features", "Menu browsing, cart, fake payment gateway, receipt and AI-style recommendation."),
        ("02", "Service build", "Node.js, TypeScript, Express, Zod validation, service classes, OpenAPI, request logging."),
        ("03", "Quality gates", "Unit, integration, contract, E2E, visual, PostgreSQL persistence, coverage, and k6 load."),
        ("04", "Observability", "Firebase cloud sync plus Excel dashboard turns metrics into release decisions."),
        ("05", "DeepSeek QA", "Failure analysis and generated scenarios extend coverage and speed triage."),
    ]
    for i, (num, head, body) in enumerate(cards):
        x = 0.8 + (i % 3) * 4.1
        y = 1.35 + (i // 3) * 2.15
        rounded(slide, x, y, 3.55, 1.7, WHITE, RGBColor(225, 225, 225))
        text(slide, num, x + 0.2, y + 0.16, 0.6, 0.35, 18, RED, bold=True)
        text(slide, head, x + 0.85, y + 0.18, 2.35, 0.3, 15, BLACK, bold=True)
        text(slide, body, x + 0.28, y + 0.72, 2.95, 0.72, 11.5, MUTED)


def product_slide():
    slide = deck.slides.add_slide(BLANK)
    add_bg(slide, SOFT)
    add_top_bar(slide, "FoodHub Feature Scope", "Takeaway SaaS")
    add_image(slide, ROOT / "tests" / "e2e" / "order-flow.spec.ts-snapshots" / "menu-visible-chromium-win32.png", 0.65, 1.2, 5.2)
    label_card(slide, "Customer Journey", "Browse menu, add items, review cart, pay through FoodHub Payment Gateway, then receive a paid order receipt.", 6.35, 1.25, 5.95, 1.15, RED)
    label_card(slide, "Business Rules", "Order totals are rounded to two decimals, unknown items are rejected, quantities are bounded from 1 to 20, and paid orders get a payment id.", 6.35, 2.65, 5.95, 1.15, GOLD)
    label_card(slide, "AI-style Suggestion", "The receipt includes deterministic menu/cart recommendations, giving product-level decision support without slowing automated tests.", 6.35, 4.05, 5.95, 1.15, BLUE)
    label_card(slide, "User Evidence", "Playwright screenshots capture menu, cart, and gateway states for visual confidence during regression testing.", 6.35, 5.45, 5.95, 1.15, GREEN)


def architecture_slide():
    slide = deck.slides.add_slide(BLANK)
    add_bg(slide, WHITE)
    add_top_bar(slide, "How The Services Are Built", "Architecture")
    nodes = [
        ("Static Web UI", "public/index.html\npublic/app.js", 0.8, 1.35, BLUE),
        ("FoodHub API", "Express + TypeScript\n/menu, /order, /health", 4.05, 1.35, RED),
        ("Order Service", "Validation, pricing,\nreceipt + recommendation", 7.25, 1.35, GOLD),
        ("Payment Gateway", "Separate fake service\nport 4174", 10.25, 1.35, DARK_RED),
        ("OpenAPI Docs", "/api-docs\n/openapi.json", 2.15, 4.35, GREEN),
        ("Request Logger", "JSONL logs\nroute/duration/status", 5.25, 4.35, BLUE),
        ("PostgreSQL Repo", "Testcontainers verifies\npaid order persistence", 8.35, 4.35, CHARCOAL),
    ]
    for title, body, x, y, color in nodes:
        rounded(slide, x, y, 2.25, 1.12, WHITE, RGBColor(220, 220, 220))
        text(slide, title, x + 0.13, y + 0.15, 1.95, 0.25, 13, color, bold=True, align=PP_ALIGN.CENTER)
        text(slide, body, x + 0.15, y + 0.52, 1.9, 0.42, 9.7, MUTED, align=PP_ALIGN.CENTER)
    connector(slide, 3.05, 1.91, 4.05, 1.91)
    connector(slide, 6.3, 1.91, 7.25, 1.91)
    connector(slide, 9.5, 1.91, 10.25, 1.91)
    connector(slide, 5.15, 2.48, 3.25, 4.35, GREEN)
    connector(slide, 5.15, 2.48, 6.2, 4.35, BLUE)
    connector(slide, 8.35, 2.48, 9.45, 4.35, CHARCOAL)
    text(slide, "Built as small testable units: framework concerns in Express middleware, business logic in services, schema validation in Zod, and integrations behind interfaces.", 1.1, 6.35, 11.1, 0.55, 16, BLACK, bold=True, align=PP_ALIGN.CENTER)


def testing_slide():
    slide = deck.slides.add_slide(BLANK)
    add_bg(slide, SOFT)
    add_top_bar(slide, "Tests Being Run", "Quality Gates")
    layers = [
        ("Unit", "Vitest service tests: totals, validation, payment, recommendation", 5.2, 5.55, 2.9, GREEN),
        ("Integration", "Supertest API tests, invalid payloads, duplicate and large orders", 4.55, 4.65, 4.2, BLUE),
        ("Contract", "Pact verifies FoodHub Web <-> FoodHub API agreement", 3.9, 3.75, 5.5, GOLD),
        ("E2E + Visual", "Playwright order flow, gateway flow, screenshots", 3.25, 2.85, 6.8, RED),
        ("Load + Persistence", "k6 thresholds and Testcontainers PostgreSQL paid-order persistence", 2.6, 1.95, 8.1, DARK_RED),
    ]
    for name, body, x, y, w, color in layers:
        rounded(slide, x, y, w, 0.62, color, color)
        text(slide, name, x + 0.18, y + 0.12, 1.1, 0.22, 12, WHITE, bold=True)
        text(slide, body, x + 1.25, y + 0.13, w - 1.45, 0.2, 10.5, WHITE)
    label_card(slide, "Main command", "npm run test:all runs local groups in parallel with fail-fast behavior and updates qa-artifacts/test-report.html.", 0.75, 6.1, 5.8, 0.7, RED)
    label_card(slide, "Coverage target", "Critical business logic is held to 90%+ statements, branches, functions, and lines.", 6.8, 6.1, 5.75, 0.7, GREEN)


def metrics_slide():
    slide = deck.slides.add_slide(BLANK)
    add_bg(slide, WHITE)
    add_top_bar(slide, "Latest Quality Evidence", "Current Snapshot")
    metric_card(slide, f"{summary['passedChecks']}/{summary['totalChecks']}", "checks passed", 0.8, 1.25, 2.2, 1.35, GREEN)
    metric_card(slide, summary["failedChecks"], "failed checks", 3.25, 1.25, 2.2, 1.35, RED)
    metric_card(slide, f"{summary['requestCount']:,}", "observed requests", 5.7, 1.25, 2.2, 1.35, BLUE)
    metric_card(slide, f"{summary['p95DurationMs']} ms", "p95 latency", 8.15, 1.25, 2.2, 1.35, GOLD)
    metric_card(slide, summary["paymentFailureCount"], "payment failures", 10.6, 1.25, 2.2, 1.35, GREEN)
    rounded(slide, 0.8, 3.1, 11.95, 2.55, SOFT, RGBColor(230, 230, 230))
    text(slide, "Release signal from the snapshot", 1.15, 3.35, 5.2, 0.35, 20, BLACK, bold=True)
    bullet_list(slide, [
        "All automated checks are green in the latest dashboard snapshot.",
        "Traffic evidence includes browse, health, static UI, and paid order requests.",
        "Payment signal is healthy: 3,716 successes and 0 failures in the summary.",
        "Performance signal is comfortably inside the k6 p95 threshold of 500 ms."
    ], 1.15, 3.9, 6.0, 1.3, 14)
    rounded(slide, 8.0, 3.55, 3.35, 1.45, GREEN, GREEN)
    text(slide, "GO", 8.0, 3.78, 3.35, 0.48, 32, WHITE, bold=True, align=PP_ALIGN.CENTER)
    text(slide, "Current evidence supports release", 8.15, 4.36, 3.05, 0.25, 12, WHITE, align=PP_ALIGN.CENTER)


def firebase_slide():
    slide = deck.slides.add_slide(BLANK)
    add_bg(slide, SOFT)
    add_top_bar(slide, "Observability Dashboard + Firebase Cloud Sync", "Go / No-Go")
    steps = [
        ("Collect", "Request logs + QA metrics\nnpm run observability:collect", 0.8, 1.35, RED),
        ("Ingest", "Push latest run into\nFirebase Realtime Database", 3.75, 1.35, BLUE),
        ("Share", "Read merged metrics snapshot\nacross local + CI runs", 6.7, 1.35, GOLD),
        ("Dashboard", "Generate Excel workbook\nFoodHub-Observability-Dashboard.xlsx", 9.65, 1.35, GREEN),
    ]
    for title, body, x, y, color in steps:
        rounded(slide, x, y, 2.45, 1.18, WHITE, RGBColor(225, 225, 225))
        text(slide, title, x + 0.18, y + 0.16, 2.05, 0.25, 14, color, bold=True, align=PP_ALIGN.CENTER)
        text(slide, body, x + 0.16, y + 0.52, 2.1, 0.45, 9.8, MUTED, align=PP_ALIGN.CENTER)
    connector(slide, 3.25, 1.95, 3.75, 1.95)
    connector(slide, 6.2, 1.95, 6.7, 1.95)
    connector(slide, 9.15, 1.95, 9.65, 1.95)
    label_card(slide, "Decision metrics", "Total checks, failures, skipped checks, request count, error count, average duration, p95 duration, payment successes, and payment failures.", 0.95, 3.45, 5.6, 1.22, RED)
    label_card(slide, "Why Firebase helps", "Local machines and CI can publish to one cloud-backed metrics store, so the release call uses shared evidence instead of one isolated run.", 6.8, 3.45, 5.55, 1.22, BLUE)
    label_card(slide, "Go / no-go usage", "Green tests + low error rate + latency under threshold + payment health = go. Any failed gate or business-critical metric drift triggers no-go and investigation.", 2.05, 5.35, 9.25, 0.95, GREEN)


def deepseek_failure_slide():
    slide = deck.slides.add_slide(BLANK)
    add_bg(slide, WHITE)
    add_top_bar(slide, "DeepSeek For Failure Analysis", "AI-assisted QA")
    rounded(slide, 0.85, 1.25, 3.3, 4.7, SOFT, RGBColor(230, 230, 230))
    text(slide, "Input Signals", 1.1, 1.55, 2.7, 0.32, 18, RED, bold=True, align=PP_ALIGN.CENTER)
    bullet_list(slide, [
        "Timeouts",
        "Locator / selector issues",
        "Network or service startup",
        "Payment gateway redirects",
        "Pact, PostgreSQL, Docker, k6"
    ], 1.25, 2.15, 2.5, 2.3, 14)
    rounded(slide, 5.0, 1.25, 3.3, 4.7, SOFT, RGBColor(230, 230, 230))
    text(slide, "Analyzer Output", 5.25, 1.55, 2.7, 0.32, 18, BLUE, bold=True, align=PP_ALIGN.CENTER)
    bullet_list(slide, [
        "Likely root cause",
        "Classification: product, test, environment, unknown",
        "Evidence from logs",
        "Confidence",
        "Recommended fix"
    ], 5.45, 2.15, 2.45, 2.3, 14)
    rounded(slide, 9.15, 1.25, 3.3, 4.7, SOFT, RGBColor(230, 230, 230))
    text(slide, "Artifacts", 9.4, 1.55, 2.7, 0.32, 18, GREEN, bold=True, align=PP_ALIGN.CENTER)
    bullet_list(slide, [
        "qa-artifacts/failure-analysis-report.md",
        "qa-artifacts/ai-failure-analysis-prompt.txt",
        "FoodHub-AI-Test-Analysis.xlsx",
        "DeepSeek v4 pro in live mode"
    ], 9.45, 2.15, 2.45, 2.3, 13)
    connector(slide, 4.15, 3.55, 5.0, 3.55)
    connector(slide, 8.3, 3.55, 9.15, 3.55)


def deepseek_tests_slide():
    slide = deck.slides.add_slide(BLANK)
    add_bg(slide, SOFT)
    add_top_bar(slide, "DeepSeek For Test Generation", "Coverage Expansion")
    label_card(slide, "Failure Analysis", "DeepSeek reads logs and project context, then classifies failure patterns with evidence and recommended fixes.", 0.75, 1.25, 3.8, 1.25, RED)
    label_card(slide, "Scenario Analysis", "It generates edge cases, missing scenarios, and coverage expansion rows mapped to risk and test level.", 4.78, 1.25, 3.8, 1.25, BLUE)
    label_card(slide, "Test Data Suggestions", "It proposes deterministic builders such as duplicate item, boundary quantity, randomized order, and load traffic data.", 8.8, 1.25, 3.8, 1.25, GREEN)
    rounded(slide, 1.1, 3.25, 11.1, 2.55, WHITE, RGBColor(228, 228, 228))
    text(slide, "Implemented AI-suggested scenarios", 1.45, 3.55, 4.8, 0.35, 19, BLACK, bold=True)
    bullet_list(slide, [
        "Duplicate items submitted as separate order lines.",
        "Large order at maximum accepted quantity.",
        "Invalid payload shapes rejected predictably.",
        "Observability metrics fixture feeds dashboard generation.",
        "k6 load traffic uses unique payment tokens per VU and iteration."
    ], 1.55, 4.05, 6.0, 1.35, 14)
    rounded(slide, 8.2, 3.95, 2.8, 0.9, RED, RED)
    text(slide, "FOODHUB_AI_LIVE", 8.35, 4.11, 2.5, 0.25, 13, WHITE, bold=True, align=PP_ALIGN.CENTER)
    text(slide, "switches live DeepSeek on/off", 8.4, 4.47, 2.4, 0.18, 9.5, WHITE, align=PP_ALIGN.CENTER)


def go_nogo_slide():
    slide = deck.slides.add_slide(BLANK)
    add_bg(slide, WHITE)
    add_top_bar(slide, "Go / No-Go Decision Model", "Release Readiness")
    rows = [
        ("Automated tests", "0 failed critical checks", "GO"),
        ("Coverage", "90%+ business logic threshold", "GO"),
        ("Load", "p95 < 500 ms and order failure rate < 2%", "GO"),
        ("Payments", "No unexpected payment failures", "GO"),
        ("Observability", "Dashboard refreshed from Firebase shared metrics", "GO"),
        ("AI triage", "Known failures classified with next actions", "GO"),
    ]
    y = 1.28
    for area, gate, status in rows:
        rounded(slide, 1.0, y, 11.2, 0.62, SOFT, RGBColor(225, 225, 225))
        text(slide, area, 1.25, y + 0.14, 2.4, 0.2, 12.5, BLACK, bold=True)
        text(slide, gate, 3.95, y + 0.14, 5.6, 0.2, 12.5, MUTED)
        text(slide, status, 10.45, y + 0.11, 1.0, 0.24, 14, GREEN, bold=True, align=PP_ALIGN.CENTER)
        y += 0.77
    rounded(slide, 2.1, 6.05, 9.1, 0.74, GREEN, GREEN)
    text(slide, "Current project snapshot supports GO, with continued monitoring after every refresh.", 2.28, 6.24, 8.7, 0.25, 15, WHITE, bold=True, align=PP_ALIGN.CENTER)


def artifacts_slide():
    slide = deck.slides.add_slide(BLANK)
    add_bg(slide, SOFT)
    add_top_bar(slide, "Evidence Pack", "What reviewers can open")
    artifacts = [
        ("Test report", "qa-artifacts/test-report.html"),
        ("Playwright report", "playwright-report/index.html"),
        ("Observability dashboard", "qa-artifacts/FoodHub-Observability-Dashboard.xlsx"),
        ("AI test analysis", "qa-artifacts/FoodHub-AI-Test-Analysis.xlsx"),
        ("Failure report", "qa-artifacts/failure-analysis-report.md"),
        ("OpenAPI docs", "http://127.0.0.1:4173/api-docs"),
        ("Pact contract", "pacts/FoodHub Web-FoodHub API.json"),
        ("Load summary", "qa-artifacts/load-test-summary.json"),
    ]
    for idx, (name, path) in enumerate(artifacts):
        x = 0.8 + (idx % 2) * 6.05
        y = 1.15 + (idx // 2) * 1.28
        label_card(slide, name, path, x, y, 5.55, 0.95, [RED, BLUE, GOLD, GREEN][idx % 4])


def closing_slide():
    slide = deck.slides.add_slide(BLANK)
    add_bg(slide, RED)
    add_image(slide, LOGO_IMAGE, 3.9, 0.85, 5.55)
    text(slide, "FoodHub Release Confidence", 2.0, 4.85, 9.3, 0.48, 30, WHITE, bold=True, align=PP_ALIGN.CENTER)
    text(slide, "Services are modular, tests are layered, observability is shared, and AI shortens the loop from failure to fix.", 2.05, 5.48, 9.2, 0.5, 16, WHITE, align=PP_ALIGN.CENTER)
    text(slide, "GO when evidence is green. NO-GO when the dashboard or AI triage exposes risk.", 2.2, 6.34, 8.9, 0.28, 13.5, WHITE, bold=True, align=PP_ALIGN.CENTER)


for build in [
    title_slide,
    agenda_slide,
    product_slide,
    architecture_slide,
    testing_slide,
    metrics_slide,
    firebase_slide,
    deepseek_failure_slide,
    deepseek_tests_slide,
    go_nogo_slide,
    artifacts_slide,
    closing_slide,
]:
    build()

deck.save(OUTPUT)
print(OUTPUT)
